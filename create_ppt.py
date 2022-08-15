import os
import glob
import shutil
import datetime
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Cm, Pt
from PIL import Image
from utils import get_initial_data

# Logger setting
from logging import getLogger, FileHandler, DEBUG
logger = getLogger(__name__)
today = datetime.datetime.now()
os.makedirs('./log', exist_ok=True)
handler = FileHandler(f'log/{today.strftime("%Y-%m-%d")}_result.log', mode='a')
handler.setLevel(DEBUG)
logger.setLevel(DEBUG)
logger.addHandler(handler)
logger.propagate = False


### functions ###
def add_picture(slide, img_file, slide_width, slide_height):
    slide_aspect_ratio = slide_width / slide_height
    img_center_x = slide_width / 2

    #画像サイズを取得してアスペクト比を得る
    im = Image.open(img_file)
    im_width, im_height = im.size
    aspect_ratio = im_width/im_height

    #スライドと画像のアスペクト比に応じて処理を分岐
    #画像のほうが横長だったら横めいっぱいに広げる
    if aspect_ratio > slide_aspect_ratio:
      img_display_width = slide_width
      img_display_height = img_display_width / aspect_ratio
    else: #画像のほうが縦長だったら縦めいっぱいに広げる
      img_display_height = slide_height - Cm(4)
      img_display_width = img_display_height * aspect_ratio
    #センタリングする場合の画像の左上座標を計算
    left = img_center_x - img_display_width / 2
    top = Cm(3)

    #画像をスライドに追加
    if aspect_ratio > slide_aspect_ratio:
      slide.shapes.add_picture(img_file, left, top, width = img_display_width)  
    else:
      slide.shapes.add_picture(img_file, left, top, height = img_display_height)  


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def set_cell_border(cell, border_color="000000", border_width='30000'):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Left Cell Border
    lnL = SubElement(tcPr, 'a:lnL', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    lnL_solidFill = SubElement(lnL, 'a:solidFill')
    lnL_srgbClr = SubElement(lnL_solidFill, 'a:srgbClr', val=border_color)
    lnL_prstDash = SubElement(lnL, 'a:prstDash', val='solid')
    lnL_round_ = SubElement(lnL, 'a:round')
    lnL_headEnd = SubElement(lnL, 'a:headEnd', type='none', w='med', len='med')
    lnL_tailEnd = SubElement(lnL, 'a:tailEnd', type='none', w='med', len='med')

    # Right Cell Border
    lnR = SubElement(tcPr, 'a:lnR', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    lnR_solidFill = SubElement(lnR, 'a:solidFill')
    lnR_srgbClr = SubElement(lnR_solidFill, 'a:srgbClr', val=border_color)
    lnR_prstDash = SubElement(lnR, 'a:prstDash', val='solid')
    lnR_round_ = SubElement(lnR, 'a:round')
    lnR_headEnd = SubElement(lnR, 'a:headEnd', type='none', w='med', len='med')
    lnR_tailEnd = SubElement(lnR, 'a:tailEnd', type='none', w='med', len='med')

    # Top Cell Border
    lnT = SubElement(tcPr, 'a:lnT', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    lnT_solidFill = SubElement(lnT, 'a:solidFill')
    lnT_srgbClr = SubElement(lnT_solidFill, 'a:srgbClr', val=border_color)
    lnT_prstDash = SubElement(lnT, 'a:prstDash', val='solid')
    lnT_round_ = SubElement(lnT, 'a:round')
    lnT_headEnd = SubElement(lnT, 'a:headEnd', type='none', w='med', len='med')
    lnT_tailEnd = SubElement(lnT, 'a:tailEnd', type='none', w='med', len='med')

    # Bottom Cell Border
    lnB = SubElement(tcPr, 'a:lnB', w=border_width, cap='flat', cmpd='sng', algn='ctr')
    lnB_solidFill = SubElement(lnB, 'a:solidFill')
    lnB_srgbClr = SubElement(lnB_solidFill, 'a:srgbClr', val=border_color)
    lnB_prstDash = SubElement(lnB, 'a:prstDash', val='solid')
    lnB_round_ = SubElement(lnB, 'a:round')
    lnB_headEnd = SubElement(lnB, 'a:headEnd', type='none', w='med', len='med')
    lnB_tailEnd = SubElement(lnB, 'a:tailEnd', type='none', w='med', len='med')


def get_article_data(sheet):
    articles = sheet.get_all_values()
    articles.pop(0)
    articles.pop(0)
    for article in articles:
        if article[7] == '○':
            yield {
                'no': article[0],
                'media': article[1],
                'date': datetime.datetime.strptime(article[2], "%Y-%m-%d"),
                'url': article[3]
            }


def create_ppt(path, data, title):
    ppt = Presentation()
    ##スライドのサイズを指定
    ppt.slide_width = Cm(19.05)
    ppt.slide_height = Cm(27.52)

    for d in data:
        ##追加するスライドを選択
        sld0 = ppt.slides.add_slide(ppt.slide_layouts[6])

        rows = 2
        cols = 4
        table_shape = sld0.shapes.add_table(rows, cols,
                                       Cm(0.62), Cm(0.5), Cm(17.81), Cm(1.98))
        table = table_shape.table
        table.columns[0].width = Cm(2.5)
        table.columns[1].width = Cm(6.4)
        table.columns[2].width = Cm(2.5)
        table.columns[3].width = Cm(6.41)
        table.cell(1, 1).merge(table.cell(1, 3))

        values = [
            {'row':0, 'col':0, 'text':'掲載日', 'type': 'header'},
            {'row':0, 'col':1, 'text':d['date'].strftime("%-m月%-d日"), 'type': 'date'},
            {'row':0, 'col':2, 'text':'媒体名', 'type': 'header'},
            {'row':0, 'col':3, 'text':d['media'], 'type': 'media'},
            {'row':1, 'col':0, 'text':'URL', 'type': 'header'},
            {'row':1, 'col':1, 'text':d['url'], 'type': 'url'}
        ]

        for v in values:
            cell = table.cell(v['row'], v['col'])
            set_cell_border(cell)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text = v['text']
            pg = cell.text_frame.paragraphs[0]
            pg.font.name = 'Meiryo UI'
            pg.font.size = Pt(14)
            cell.fill.solid()
            if v['type'] == 'header':
                cell.fill.fore_color.rgb = RGBColor(255, 53, 28)
                pg.font.color.rgb = RGBColor(255, 255, 255)
                pg.font.bold = True
                #pg.alignment = PP_ALIGN.CENTER
            elif v['type'] == 'url':
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                pg.font.color.rgb = RGBColor(0, 0, 255)
                pg.font.underline = True
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                pg.font.color.rgb = RGBColor(0, 0, 0)
        
        add_picture(sld0, f"./screenshots/{title}/{d['no']}.png", ppt.slide_width, ppt.slide_height)

    ppt.save(path)

### main_script ###
if __name__ == '__main__':
    logger.info('\n----- Start getting initial data -----\n')
    projects = get_initial_data()
    logger.info(f'Projects: size: {len(projects)}')

    logger.info('\n----- Start create powerpoint -----\n')
    os.makedirs('./ppt', exist_ok=True)
    try:
        for index, p in enumerate(projects):
            logger.info(f"{index}: {p['sheet'].title}")

            path = f"./ppt/【掲載記事クリッピング】{p['sheet'].title}.ppt" 
            is_file = os.path.isfile(path)
            if is_file:
                os.remove(path)

            data = get_article_data(p['sheet'])
            create_ppt(path, data, p['sheet'].title)
            files = glob.glob(f"./screenshots/{p['sheet'].title}/*")
            for f in files:
                print(f)
            exit(0)

    except Exception as err:
        logger.debug(f'{err}')
        exit(1)